import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Humanized Dark Defense Palette
BG_COLOR = RGBColor(11, 15, 25)          # Deep Navy #0B0F19
PANEL_COLOR = RGBColor(22, 30, 46)       # Slate Card #161E2E
PANEL_BORDER = RGBColor(45, 58, 80)      # Slate Border #2D3A50
GOLD_ACCENT = RGBColor(245, 158, 11)     # Warm Saffron Gold #F59E0B
GREEN_ACCENT = RGBColor(16, 185, 129)    # Emerald Green #10B981
BLUE_ACCENT = RGBColor(59, 130, 246)     # Sky Blue #3B82F6
PURPLE_ACCENT = RGBColor(168, 85, 247)   # Purple #A855F7
RED_ACCENT = RGBColor(239, 68, 68)       # Alert Crimson #EF4444
WHITE = RGBColor(255, 255, 255)
MUTED_TEXT = RGBColor(148, 163, 184)     # Light Slate #94A3B8

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, tag, title, subtitle=None):
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag.upper()
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = GOLD_ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = MUTED_TEXT

def add_card(slide, left, top, width, height, title, body_bullets, border_color=None, bg_color=PANEL_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = PANEL_BORDER
        shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        p_title.alignment = PP_ALIGN.LEFT

    for i, bullet in enumerate(body_bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_TEXT if not bullet.startswith("✓") and not bullet.startswith("★") and not bullet.startswith("•") else WHITE
        p.space_before = Pt(5)
        p.alignment = PP_ALIGN.LEFT

slide_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 1: HUMANIZED TITLE
# ============================================================
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1)

badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(4.8), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = GOLD_ACCENT
badge.line.fill.background()
tf_badge = badge.text_frame
p_b = tf_badge.paragraphs[0]
p_b.text = "HUMAN-CENTERED FINANCIAL DEFENSE"
p_b.font.size = Pt(11)
p_b.font.bold = True
p_b.font.color.rgb = RGBColor(11, 15, 25)
p_b.alignment = PP_ALIGN.CENTER

title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p1 = tf.paragraphs[0]
p1.text = "AI DEFENSE LAB"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = WHITE

p2 = tf.add_paragraph()
p2.text = "Protecting People, Payments & Networks with the 3-Layer Defense Triad"
p2.font.size = Pt(22)
p2.font.color.rgb = GOLD_ACCENT
p2.font.bold = True

p3 = tf.add_paragraph()
p3.text = "\"Because modern fraud doesn't hack computers — it hacks human psychology.\""
p3.font.size = Pt(14)
p3.font.italic = True
p3.font.color.rgb = MUTED_TEXT

add_card(slide1, Inches(0.8), Inches(4.3), Inches(3.6), Inches(2.2), "1. Protect the Human", [
    "• Captures live keystroke rhythm & stress",
    "• Detects active phone calls & coercion",
    "• Intercepts Digital Arrest scams early"
], GREEN_ACCENT)

add_card(slide1, Inches(4.8), Inches(4.3), Inches(3.6), Inches(2.2), "2. Secure the Transaction", [
    "• Server-side anti-tamper feature store",
    "• Evaluates 1,000 transactions in 32ms",
    "• TreeSHAP exact mathematical proof"
], BLUE_ACCENT)

add_card(slide1, Inches(8.8), Inches(4.3), Inches(3.6), Inches(2.2), "3. Trace the Network", [
    "• Uncovers cyclic money laundering loops",
    "• Detects 10-node smurfing fan-outs",
    "• Freezes entire mule rings in one action"
], PURPLE_ACCENT)

footer_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4))
p_foot = footer_box.text_frame.paragraphs[0]
p_foot.text = "Team: IRON  |  Live Prototype: https://defense-six.vercel.app/  |  GitHub: https://github.com/Abhay-s-8/defense"
p_foot.font.size = Pt(12)
p_foot.font.color.rgb = GOLD_ACCENT

# ============================================================
# SLIDE 2: THE REAL HUMAN TRAGEDY (DIGITAL ARREST SCAMS)
# ============================================================
slide2 = prs.slides.add_slide(slide_layout)
set_slide_background(slide2)
add_header(slide2, "The Real-World Crisis", "The ₹50 Lakh Phone Call: Why Traditional Bank AI Fails", "When victims are manipulated into sending money themselves, traditional algorithms see zero anomalies.")

add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "A Real-World Scenario", [
    "1. The Setup:",
    "   • A 62-year-old retired citizen receives a fake CBI/Police video call.",
    "   • Scammers claim their Aadhaar is linked to money laundering.",
    "   • Under extreme fear, they are ordered to transfer ₹10 Lakhs immediately.",
    "",
    "2. The Transaction:",
    "   • The victim opens their own banking app on their regular phone.",
    "   • Enters their genuine PIN, biometric fingerprint, and OTP.",
    "",
    "3. Why Traditional AI Approves It:",
    "   • Device ID is genuine, IP address is legitimate home Wi-Fi.",
    "   • Zero failed password attempts.",
    "   • Result: Life savings vanish in seconds without an alarm."
], RED_ACCENT)

add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "The 3 Blind Spots in Modern Banking", [
    "1. The Coerced Human (Layer 1 Blind Spot):",
    "   • Traditional models don't look at user interaction stress, long hesitation pauses, or active calls during payment.",
    "",
    "2. The Distributed Mule Web (Layer 3 Blind Spot):",
    "   • Attackers break stolen money into 200 micro-payments of ₹2,000 across sleeper accounts. Row-by-row models miss the network.",
    "",
    "3. The Static Rule Shield:",
    "   • Once deployed, banking models stay frozen while cybercriminals adapt low-and-slow stealth tactics.",
    "",
    "★ Our Mission: Build a humanized, self-healing 3-layer defense."
], GOLD_ACCENT)

# ============================================================
# SLIDE 3: THE 3-LAYER DEFENSE TRIAD PHILOSOPHY
# ============================================================
slide3 = prs.slides.add_slide(slide_layout)
set_slide_background(slide3)
add_header(slide3, "Our Solution", "The 3-Layer Safety Net: From Keystroke to Settlement", "Protecting transactions across the three distinct spaces where financial crime actually happens.")

add_card(slide3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Layer 1: The Human Space", [
    "Behavioral Biometrics",
    "Focus: Coercion & Interaction Rhythm",
    "",
    "• Captures live keystroke dwell & flight times",
    "• Flags long hesitation pauses (>4.5s)",
    "• Detects ongoing voice calls during input",
    "• Stops Digital Arrest scams at checkout",
    "",
    "✓ Fusion Weight: 20%"
], GREEN_ACCENT)

add_card(slide3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Layer 2: The Data Space", [
    "Fortified XGBoost V4 Engine",
    "Focus: Speed, Velocity & Math",
    "",
    "• Server-side real-time feature store",
    "• Anti-injection: zero client metadata trust",
    "• 1,000 transactions scored in 32ms",
    "• TreeSHAP exact mathematical weights",
    "",
    "✓ Fusion Weight: 50%"
], BLUE_ACCENT)

add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Layer 3: The Network Space", [
    "GNN Mule Ring Analyzer",
    "Focus: Account Graphs & Smurfing",
    "",
    "• Directed graph tracking money flows",
    "• Cyclic fund routing (A → B → C → A)",
    "• Fan-out smurfing & mule aggregation",
    "• Freezes entire syndicates at once",
    "",
    "✓ Fusion Weight: 30%"
], PURPLE_ACCENT)

# ============================================================
# SLIDE 4: LAYER 1 IN ACTION (BIOMETRICS & COERCION)
# ============================================================
slide4 = prs.slides.add_slide(slide_layout)
set_slide_background(slide4)
add_header(slide4, "Layer 1 In Action", "Protecting the Coerced User: Keystroke & Call Telemetry", "How we catch psychological distress and automated scripts before a payment is dispatched.")

add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "What the Engine Senses in Real Time", [
    "1. Live Typing Rhythm (Dwell & Flight Times):",
    "   • Genuine humans type with natural cadence (~90-160ms per key).",
    "   • Bots type with robotic 0ms variance or inhuman speed (<25ms).",
    "",
    "2. The Psychological Distress Signature:",
    "   • Extreme hesitation pause on the 'Pay' button (>4.5 seconds).",
    "   • Elevated backspace count (anxious corrections).",
    "   • Active voice call active during checkout (Vishing marker).",
    "",
    "3. Remote Screen Overlays:",
    "   • Detects AnyDesk/TeamViewer background hooks."
], GREEN_ACCENT)

add_card(slide4, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "The Life-Saving Safety Intervention", [
    "What Happens When Coercion is Sensed:",
    "",
    "Step 1: User types payment while pressured on a scam call.",
    "Step 2: Layer 1 detects active call + 6.8s submit hesitation.",
    "Step 3: Unified Engine triggers policy action:",
    "",
    "★ DIRECTIVE: COERCION_SAFETY_INTERVENTION",
    "★ ACTION: Pops an unclosable scam warning modal & enforces a 15-minute cooling-off delay.",
    "",
    "✓ Result: Breaks the scammer's psychological grip; citizen calls family/police before money leaves."
], GOLD_ACCENT)

# ============================================================
# SLIDE 5: LAYER 2 IN ACTION (XGBOOST & ANTI-INJECTION)
# ============================================================
slide5 = prs.slides.add_slide(slide_layout)
set_slide_background(slide5)
add_header(slide5, "Layer 2 In Action", "Sub-Millisecond Fraud Math: Zero Client Trust", "How we score 1,000 transactions in 32ms and eliminate client-side feature tampering.")

add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Anti-Tamper Real-Time Feature Store", [
    "The Fatal Flaw in Many Fraud APIs:",
    "• Naive systems trust client payloads (e.g. client says velocity is 1).",
    "• Hackers edit the JSON payload to fool the model.",
    "",
    "Our Fix (RealTimeFeatureStore):",
    "• Accepts raw minimal payment (₹ Amount, Merchant, Card).",
    "• Server computes rolling 5-minute velocity independently.",
    "• Server derives historical 30-day spending averages.",
    "• Threat-intel IP reputation verified server-side.",
    "",
    "✓ Zero Client Trust & Tamper-Proof Ingestion."
], BLUE_ACCENT)

add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "32ms Vectorized Speed & TreeSHAP Math", [
    "Vectorized Batch Inference:",
    "• Single-pass Pandas matrix manipulation.",
    "• Evaluates 1,000 transactions in 32 milliseconds (0.032s).",
    "• Ready for high-volume gateways like UPI / NPCI.",
    "",
    "Exact TreeSHAP Mathematical Proof:",
    "• Extracts exact booster tree split weights (`pred_contribs=True`).",
    "• Identifies positive risk drivers: `[+] Velocity (+3.44)`",
    "• Identifies negative safety anchors: `[-] Device Age (-0.74)`",
    "",
    "✓ Eliminates LLM hallucinations in compliance reports."
], BLUE_ACCENT)

# ============================================================
# SLIDE 6: LAYER 3 IN ACTION (GNN MULE NETWORKS)
# ============================================================
slide6 = prs.slides.add_slide(slide_layout)
set_slide_background(slide6)
add_header(slide6, "Layer 3 In Action", "Detective Work on Graphs: Dismantling Mule Rings", "How Graph Neural Network analysis catches coordinated smurfing and cyclic money laundering.")

add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "How Graph Link Analysis Works", [
    "1. Directed Fund-Flow Topology:",
    "   • Maps real-time relationships between accounts, cards, and merchants.",
    "",
    "2. Cyclic Money Laundering Detection:",
    "   • Detects closed wash-trading loops (A → B → C → A).",
    "   • Catches syndicates recycling funds to disguise their origin.",
    "",
    "3. Smurfing Fan-Out & Aggregation:",
    "   • Flags accounts rapidly dispersing micro-transfers to 4+ sleeper mules.",
    "   • Flags destination accounts with rapid high in-degree inflows."
], PURPLE_ACCENT)

add_card(slide6, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "The Syndicate Freeze Action", [
    "Scenario: Criminal syndicate laundering ₹25 Lakhs.",
    "",
    "Step 1: Stolen funds dispersed across 5 sleeper mule accounts.",
    "Step 2: GNN engine traces directed edges and flags high degree fan-out.",
    "Step 3: Mules attempt to funnel funds back to an aggregator.",
    "Step 4: Cyclic flow algorithm flags closed laundering loop.",
    "",
    "★ DIRECTIVE: MULE_SYNDICATE_FREEZE",
    "★ ACTION: Isolates the entire sub-graph and freezes all connected mule accounts in one operation."
], PURPLE_ACCENT)

# ============================================================
# SLIDE 7: DECISION FUSION & ACTION DIRECTIVES
# ============================================================
slide7 = prs.slides.add_slide(slide_layout)
set_slide_background(slide7)
add_header(slide7, "Policy Intelligence", "Unified Decision Fusion: Frictionless Yet Impenetrable", "Balancing zero-friction everyday purchases with targeted interventions against real threats.")

add_card(slide7, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "1. Frictionless Approval", [
    "Risk Score: 0% - 34%",
    "Color: Emerald Green",
    "",
    "• Genuine biometrics & normal cadence",
    "• Expected spend within 30d baseline",
    "• Clean graph node topology",
    "",
    "✓ 99% of normal users enjoy instant, seamless checkout without friction."
], GREEN_ACCENT)

add_card(slide7, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "2. Step-Up 2FA Challenge", [
    "Risk Score: 35% - 69%",
    "Color: Amber Yellow",
    "",
    "• Moderate behavioral velocity",
    "• New device or late-night timing",
    "• No coercion flags detected",
    "",
    "✓ Prompts for FIDO2 WebAuthn or Biometric Passkey verification."
], GOLD_ACCENT)

add_card(slide7, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "3. Targeted Interventions", [
    "Specialized Defense Directives:",
    "",
    "★ COERCION_SAFETY_INTERVENTION",
    "• Triggered by voice call + severe hesitation.",
    "• Introduces 15-min scam cooling delay.",
    "",
    "★ MULE_SYNDICATE_FREEZE",
    "• Triggered by cyclic laundering loop / smurfing.",
    "• Freezes entire multi-node cluster."
], RED_ACCENT)

# ============================================================
# SLIDE 8: THE SELF-HEALING ADVERSARIAL RED TEAM
# ============================================================
slide8 = prs.slides.add_slide(slide_layout)
set_slide_background(slide8)
add_header(slide8, "Adversarial Lab", "The Red Team Simulator: An AI That Attacks Itself", "Constantly stress-testing our defenses with 6 distinct attack families to eliminate zero-day blind spots.")

add_card(slide8, Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.3), "1. Account Takeover", [
    "• New device + high IP risk + location shift",
    "• Repeated password failures prior to tx"
], RED_ACCENT)

add_card(slide8, Inches(4.8), Inches(1.8), Inches(3.6), Inches(2.3), "2. Card Testing", [
    "• Rapid velocity (15-30 txs / 5 mins)",
    "• Micro-amounts (₹1 - ₹100 probing)"
], RED_ACCENT)

add_card(slide8, Inches(8.8), Inches(1.8), Inches(3.6), Inches(2.3), "3. Low & Slow Stealth", [
    "• Mimics normal daytime spend",
    "• Near-baseline amounts (0.8x - 1.2x)"
], RED_ACCENT)

add_card(slide8, Inches(0.8), Inches(4.3), Inches(3.6), Inches(2.3), "4. Rapid Mule Drain", [
    "• High velocity to newly created beneficiary",
    "• High amount deviation"
], RED_ACCENT)

add_card(slide8, Inches(4.8), Inches(4.3), Inches(3.6), Inches(2.3), "5. Digital Arrest Coercion", [
    "• Authentic device + active phone call",
    "• Severe hesitation and backspaces"
], GOLD_ACCENT)

add_card(slide8, Inches(8.8), Inches(4.3), Inches(3.6), Inches(2.3), "6. Mule Syndicates", [
    "• Multi-account cyclic fund loops",
    "• Fan-out smurfing across clusters"
], PURPLE_ACCENT)

# ============================================================
# SLIDE 9: CLOSED-LOOP RETRAINING & BALANCED REPLAY
# ============================================================
slide9 = prs.slides.add_slide(slide_layout)
set_slide_background(slide9)
add_header(slide9, "Self-Healing AI", "Closed-Loop Retraining: Learning from Every Miss", "How our model turns missed attacks into defensive armor while keeping False Alarms under 0.8%.")

add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "The Secret: Balanced Replay Buffers", [
    "The Danger of Naive Retraining:",
    "• If a model only trains on new fraud, it forgets what normal users look like (Catastrophic Forgetting).",
    "• Result: Genuine customers get blocked, causing outrage.",
    "",
    "Our Balanced Replay Solution:",
    "• Captures missed Red Team attacks.",
    "• Injects 1,200 normal benign transactions alongside them.",
    "• Retrains XGBoost with cross-entropy loss and early stopping.",
    "• Serializes fortified weights directly to disk.",
    "",
    "✓ Retains 100% normal user accuracy with zero forgetting."
], GREEN_ACCENT)

add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Adversarial Benchmark Results", [
    "Round 2 (Initial Baseline Model):",
    "• 5,000 Stealth Adversarial Attacks Launched",
    "• Detected: 582  |  Missed: 4,418",
    "• Detection Rate: 11.64%",
    "",
    "Round 3 (After Closed-Loop Retraining):",
    "• 5,000 Held-Out Adversarial Attacks Launched",
    "• Detected: 4,874  |  Missed: 126",
    "• Detection Rate: 97.48%",
    "",
    "★ False Negatives Slashed by 97.15%",
    "★ False Positive Rate Held at <0.8%"
], GOLD_ACCENT)

# ============================================================
# SLIDE 10: GROUNDED GENAI EXPLAINABILITY
# ============================================================
slide10 = prs.slides.add_slide(slide_layout)
set_slide_background(slide10)
add_header(slide10, "Explainable AI (XAI)", "Grounded GenAI: Explanations Humans & Regulators Trust", "Translating exact mathematical TreeSHAP feature weights into clear, compliance-ready English.")

add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Why Mathematical Grounding Matters", [
    "The Risk with Ungrounded LLMs:",
    "• Standard GenAI hallucinates explanations based on superficial guesses.",
    "• Fails strict banking regulations (RBI, GDPR, OCC).",
    "",
    "Our TreeSHAP Grounded Approach:",
    "• Booster tree split margins extracted mathematically.",
    "• Top positive risk drivers: `[+] Velocity (+2.84)`, `[+] IP Risk (+2.70)`",
    "• Top negative safety anchors: `[-] Account Age (-0.13)`",
    "• LLM prompt strictly constrained to mathematical weights.",
    "",
    "✓ 100% Mathematically Grounded & Regulatory Compliant."
], BLUE_ACCENT)

add_card(slide10, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Sample Plain-Language Output", [
    "What the Fraud Officer & Regulator Sees:",
    "",
    "\"Transaction flagged as FRAUD (99.31% probability).",
    "The decision is mathematically driven by an abnormal surge in 5-minute velocity (+2.84 risk driver) and an elevated IP risk score (+2.70 risk driver).",
    "",
    "Key Signals Observed:",
    "• [+] Transaction Velocity 5M (+2.84 risk driver)",
    "• [+] IP Risk Score (+2.70 risk driver)",
    "• [-] Account Age Days (-0.13 safety anchor)",
    "",
    "Recommended Action: Trigger FIDO2 step-up challenge or route to Tier-2 fraud review.\""
], GOLD_ACCENT)

# ============================================================
# SLIDE 11: PRODUCTION HARDENING & RESILIENCY
# ============================================================
slide11 = prs.slides.add_slide(slide_layout)
set_slide_background(slide11)
add_header(slide11, "Engineering Rigor", "Built for the Real World: Resilient, Fast & Secure", "Engineered for 100% demo uptime, high concurrency, and seamless cloud deployment.")

add_card(slide11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Resilient Storage", [
    "Hybrid Database Engine:",
    "",
    "• Primary: MongoDB Atlas.",
    "• Fallback: In-memory store.",
    "• Zero startup crashes.",
    "• 100% uptime guarantee."
], GREEN_ACCENT)

add_card(slide11, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "High-Speed Inference", [
    "Vectorized Performance:",
    "",
    "• Vectorized Pandas batching.",
    "• 1,000 attacks scored in 32ms.",
    "• Non-blocking workers.",
    "• Sub-second Red Team runs."
], BLUE_ACCENT)

add_card(slide11, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "API Hardening", [
    "Security & Protection:",
    "",
    "• Sliding-window rate limiting.",
    "• Brute force defense on auth/predict.",
    "• Enterprise security headers.",
    "• Model artifact disk serialization."
], GOLD_ACCENT)

# ============================================================
# SLIDE 12: THE VISION & SUMMARY
# ============================================================
slide12 = prs.slides.add_slide(slide_layout)
set_slide_background(slide12)
add_header(slide12, "Summary & Vision", "The Future of Payment Defense", "From reactive static scoring to continuous, multi-dimensional adversarial defense.")

add_card(slide12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Key Breakthroughs Accomplished", [
    "1. The 3-Layer Defense Triad:",
    "   • Layer 1 (Biometrics) stops Digital Arrest & Social Engineering.",
    "   • Layer 2 (XGBoost) stops takeover & velocity anomalies in 32ms.",
    "   • Layer 3 (GNN) dismantles multi-hop mule syndicates.",
    "",
    "2. Self-Healing Adversarial Learning:",
    "   • Closed-loop retraining boosts detection from 11.6% to 97.48%.",
    "   • Balanced replay buffers hold False Positive Rate under 0.8%.",
    "",
    "3. Grounded Explainable AI:",
    "   • Exact TreeSHAP mathematical feature weighting."
], GREEN_ACCENT)

add_card(slide12, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Team IRON — Final Takeaway", [
    "\"Fraud is not one-dimensional, and defense cannot be static.\"",
    "",
    "By uniting interaction biometrics, high-speed transaction ML, and graph intelligence into one continuous learning loop, we establish a robust, self-healing payment defense perimeter.",
    "",
    "Attack. Detect. Learn. Defend.",
    "",
    "★ Team: IRON",
    "★ Live Prototype: https://defense-six.vercel.app/",
    "★ GitHub Repo: https://github.com/Abhay-s-8/defense"
], GOLD_ACCENT)

output_path = "d:/ecg/AI_Defense_Lab_Presentation.pptx"
prs.save(output_path)
print(f"Humanized presentation saved successfully to: {output_path}")

# Copy into docs folder
os.makedirs("d:/ecg/defenselab/AI-DEFENSE-LAB/docs", exist_ok=True)
shutil.copy(output_path, "d:/ecg/defenselab/AI-DEFENSE-LAB/docs/AI_Defense_Lab_Presentation.pptx")
print("Synced to repository docs folder.")
